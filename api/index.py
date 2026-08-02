from fastapi import FastAPI, Request, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import List, Optional
import os
import json
import base64
import httpx
import anthropic

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------- Rate limiting ----------
from collections import defaultdict
import time

_rate_limit = defaultdict(list)
RATE_LIMIT_WINDOW = 60
RATE_LIMIT_MAX = 10

def _client_ip(request: Request) -> str:
    forwarded = request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    if request.client:
        return request.client.host
    return "unknown"

def _check_rate_limit(ip: str):
    now = time.time()
    window = _rate_limit[ip]
    while window and window[0] < now - RATE_LIMIT_WINDOW:
        window.pop(0)
    if len(window) >= RATE_LIMIT_MAX:
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Try again in a minute.")
    window.append(now)

# ---------- Models ----------
class CardVariable(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    symbol: Optional[str] = None
    meaning: Optional[str] = None

class Card(BaseModel):
    front: str = Field(..., min_length=1)
    back: str = Field(..., min_length=1)
    type: str = Field(default="basic", pattern="^(basic|cloze|formula)$")
    formula: Optional[str] = None
    variables: Optional[List[CardVariable]] = None
    assumptions: Optional[str] = None
    commonMistakes: Optional[str] = None
    applications: Optional[str] = None

class CardBatch(BaseModel):
    summary: str = Field(..., min_length=1)
    cards: List[Card] = Field(..., min_length=1)

class GenerateResponse(BaseModel):
    cards: List[Card]
    summary: str

# ---------- Prompts & Tools ----------
SYSTEM_PROMPT = (
    "You are a flashcard generator. Extract key concepts from the user's document "
    "and return ONLY a valid JSON object matching the submit_cards tool schema. "
    "Do not wrap the JSON in markdown fences. Do not add commentary."
)

CLAUDE_MODEL = os.environ.get("CLAUDE_MODEL", "claude-3-5-sonnet-20241022")
GEMINI_MODEL = os.environ.get("GEMINI_MODEL", "gemini-1.5-flash-latest")

GENERATE_CARDS_TOOL = {
    "name": "submit_cards",
    "description": "Submit generated flashcards",
    "input_schema": {
        "type": "object",
        "properties": {
            "summary": {"type": "string", "description": "1-2 sentence summary of the document"},
            "cards": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "front": {"type": "string"},
                        "back": {"type": "string"},
                        "type": {"type": "string", "enum": ["basic", "cloze", "formula"]},
                        "formula": {"type": ["string", "null"]},
                        "variables": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "name": {"type": ["string", "null"]},
                                    "description": {"type": ["string", "null"]},
                                    "symbol": {"type": ["string", "null"]},
                                    "meaning": {"type": ["string", "null"]},
                                },
                            },
                        },
                        "assumptions": {"type": ["string", "null"]},
                        "commonMistakes": {"type": ["string", "null"]},
                        "applications": {"type": ["string", "null"]},
                    },
                    "required": ["front", "back", "type"],
                },
            },
        },
        "required": ["summary", "cards"],
    },
}

GEMINI_RESPONSE_SCHEMA = {
    "type": "object",
    "properties": {
        "summary": {"type": "string"},
        "cards": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "front": {"type": "string"},
                    "back": {"type": "string"},
                    "type": {"type": "string", "enum": ["basic", "cloze", "formula"]},
                    "formula": {"type": ["string", "null"]},
                    "variables": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": ["string", "null"]},
                                "description": {"type": ["string", "null"]},
                                "symbol": {"type": ["string", "null"]},
                                "meaning": {"type": ["string", "null"]},
                            },
                        },
                    },
                    "assumptions": {"type": ["string", "null"]},
                    "commonMistakes": {"type": ["string", "null"]},
                    "applications": {"type": ["string", "null"]},
                },
                "required": ["front", "back", "type"],
            },
        },
    },
    "required": ["summary", "cards"],
}

# ---------- Helpers ----------
def _resolve_credentials(request: Request):
    provider = request.headers.get("x-llm-provider", "claude").lower()
    api_key = request.headers.get("x-llm-api-key", "")
    if not api_key:
        raise HTTPException(status_code=401, detail="Missing X-LLM-Api-Key header")
    if provider not in ("claude", "gemini"):
        raise HTTPException(status_code=400, detail="Unsupported provider. Use 'claude' or 'gemini'.")
    return provider, api_key

def _call_claude(text: str, provider: str, api_key: str):
    client = anthropic.Anthropic(api_key=api_key)
    response = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=4096,
        system=SYSTEM_PROMPT,
        tools=[GENERATE_CARDS_TOOL],
        tool_choice={"type": "tool", "name": "submit_cards"},
        messages=[{"role": "user", "content": text}],
    )
    for block in response.content:
        if block.type == "tool_use" and block.name == "submit_cards":
            batch = CardBatch.model_validate(block.input)
            return batch.cards, batch.summary
    raise ValueError("Model did not return a submit_cards tool call")

def _call_gemini(text: str, provider: str, api_key: str):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"
    payload = {
        "system_instruction": {"parts": [{"text": SYSTEM_PROMPT}]},
        "contents": [
            {
                "role": "user",
                "parts": [{"text": text}],
            }
        ],
        "generationConfig": {
            "responseMimeType": "application/json",
            "responseSchema": GEMINI_RESPONSE_SCHEMA,
        },
    }
    with httpx.Client(timeout=120.0) as http:
        response = http.post(url, params={"key": api_key}, json=payload)
        response.raise_for_status()
        data = response.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        parsed = json.loads(text)
        batch = CardBatch.model_validate(parsed)
        return batch.cards, batch.summary

def _extract_ppt_text(content: bytes) -> str:
    try:
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n\n".join(texts)
    except Exception:
        return ""

def _call_claude_vision(base64_data: str, mime_type: str, provider: str, api_key: str):
    client = anthropic.Anthropic(api_key=api_key)

    if mime_type == "application/pdf":
        content_blocks = [
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": base64_data
                }
            },
            {
                "type": "text",
                "text": "Generate flashcards from this document."
            }
        ]
    else:
        content_blocks = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": mime_type,
                    "data": base64_data
                }
            },
            {
                "type": "text",
                "text": "Generate flashcards from this image."
            }
        ]

    response = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=4096,
        system=SYSTEM_PROMPT,
        tools=[GENERATE_CARDS_TOOL],
        tool_choice={"type": "tool", "name": "submit_cards"},
        messages=[{"role": "user", "content": content_blocks}]
    )

    for block in response.content:
        if block.type == "tool_use" and block.name == "submit_cards":
            batch = CardBatch.model_validate(block.input)
            return batch.cards, batch.summary
    raise ValueError("Model did not return a submit_cards tool call")

def _call_gemini_vision(base64_data: str, mime_type: str, provider: str, api_key: str):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"
    payload = {
        "system_instruction": {"parts": [{"text": SYSTEM_PROMPT}]},
        "contents": [{
            "role": "user",
            "parts": [
                {
                    "inline_data": {
                        "mime_type": mime_type,
                        "data": base64_data
                    }
                },
                {"text": "Generate flashcards from this document."}
            ]
        }],
        "generationConfig": {
            "responseMimeType": "application/json",
            "responseSchema": GEMINI_RESPONSE_SCHEMA
        }
    }
    with httpx.Client(timeout=120.0) as http:
        response = http.post(url, params={"key": api_key}, json=payload)
        response.raise_for_status()
        data = response.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        parsed = json.loads(text)
        batch = CardBatch.model_validate(parsed)
        return batch.cards, batch.summary

# ---------- Endpoints ----------
@app.post("/api/generate-cards", response_model=GenerateResponse)
async def generate_cards(request: Request):
    _check_rate_limit(_client_ip(request))
    provider, api_key = _resolve_credentials(request)

    body = await request.json()
    text = body.get("text", "")
    if not text or len(text.strip()) < 10:
        raise HTTPException(status_code=400, detail="Text too short or missing.")

    try:
        if provider == "claude":
            cards, summary = _call_claude(text, provider, api_key)
        else:
            cards, summary = _call_gemini(text, provider, api_key)
    except anthropic.AuthenticationError:
        raise HTTPException(status_code=401, detail="Invalid Claude API key.")
    except anthropic.RateLimitError:
        raise HTTPException(status_code=429, detail="Claude rate limit hit. Wait a moment and retry.")
    except httpx.HTTPStatusError as e:
        raise HTTPException(status_code=502, detail=f"Gemini error: {e.response.status_code}")
    except Exception:
        # Never send internal details to clients.
        raise HTTPException(
            status_code=500,
            detail="Card generation failed. Please try again."
        )

    return GenerateResponse(cards=cards, summary=summary)

@app.post("/api/generate-cards-vision", response_model=GenerateResponse)
async def generate_cards_vision(
    request: Request,
    file: UploadFile = File(...),
    deck_id: str = Form(...)
):
    _check_rate_limit(_client_ip(request))
    provider, api_key = _resolve_credentials(request)

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="No file uploaded")

    if len(content) > 20 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large. Max 20MB.")

    mime_type = file.content_type or "application/octet-stream"
    filename = file.filename or "upload"
    ext = filename.split('.')[-1].lower() if '.' in filename else ''

    # For PowerPoint: text extraction only (a .pptx is not an image)
    if ext in ('ppt', 'pptx'):
        text = _extract_ppt_text(content)
        if text and len(text.strip()) > 50:
            try:
                if provider == "claude":
                    cards, summary = _call_claude(text, provider, api_key)
                else:
                    cards, summary = _call_gemini(text, provider, api_key)
                return GenerateResponse(cards=cards, summary=summary)
            except anthropic.AuthenticationError:
                raise HTTPException(status_code=401, detail="Invalid Claude API key.")
            except anthropic.RateLimitError:
                raise HTTPException(status_code=429, detail="Claude rate limit hit. Wait a moment and retry.")
            except httpx.HTTPStatusError as e:
                raise HTTPException(status_code=502, detail=f"Gemini error: {e.response.status_code}")
            except Exception:
                raise HTTPException(
                    status_code=500,
                    detail="Card generation failed. Please try again."
                )
        raise HTTPException(
            status_code=400,
            detail="Could not extract text from this PowerPoint. Use the manual paste flow, or export slides as PDF/images."
        )

    # For images and PDFs, use vision API
    allowed_mimes = {
        'application/pdf',
        'image/jpeg',
        'image/png',
        'image/jpg',
        'image/webp',
    }
    allowed_exts = {'pdf', 'jpg', 'jpeg', 'png', 'webp'}
    if mime_type not in allowed_mimes and ext not in allowed_exts:
        raise HTTPException(
            status_code=400,
            detail="Unsupported file type. Use PDF, JPG, or PNG — or the manual paste flow."
        )

    base64_data = base64.b64encode(content).decode('utf-8')

    try:
        if provider == "claude":
            cards, summary = _call_claude_vision(base64_data, mime_type, provider, api_key)
        else:
            cards, summary = _call_gemini_vision(base64_data, mime_type, provider, api_key)
    except anthropic.AuthenticationError:
        raise HTTPException(status_code=401, detail="Invalid Claude API key.")
    except anthropic.RateLimitError:
        raise HTTPException(status_code=429, detail="Claude rate limit hit. Wait a moment and retry.")
    except httpx.HTTPStatusError as e:
        raise HTTPException(status_code=502, detail=f"Gemini error: {e.response.status_code}")
    except Exception:
        # Never send internal details to clients.
        raise HTTPException(
            status_code=500,
            detail="Card generation failed. Please try again."
        )

    return GenerateResponse(cards=cards, summary=summary)